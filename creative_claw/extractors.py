from __future__ import annotations

import csv
import json
import mimetypes
import re
from html.parser import HTMLParser
from pathlib import Path
from typing import Callable

from .models import ExtractedUnit


class _TextHTMLParser(HTMLParser):
    BLOCK_TAGS = {"p", "div", "section", "article", "li", "h1", "h2", "h3", "h4", "h5", "h6", "br", "tr"}

    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:  # noqa: ANN001
        if tag.lower() in self.BLOCK_TAGS:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in self.BLOCK_TAGS:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        self.parts.append(data)

    def text(self) -> str:
        return re.sub(r"\n{3,}", "\n\n", "".join(self.parts)).strip()


def _read_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def extract_plain(path: Path) -> list[ExtractedUnit]:
    return [ExtractedUnit(_read_text(path), {"source_type": path.suffix.lower().lstrip(".") or "text"})]


def extract_json(path: Path) -> list[ExtractedUnit]:
    payload = json.loads(_read_text(path))
    rendered = json.dumps(payload, ensure_ascii=False, indent=2)
    return [ExtractedUnit(rendered, {"source_type": "json"})]


def extract_html(path: Path) -> list[ExtractedUnit]:
    parser = _TextHTMLParser()
    parser.feed(_read_text(path))
    return [ExtractedUnit(parser.text(), {"source_type": "html"})]


def extract_docx(path: Path) -> list[ExtractedUnit]:
    from docx import Document

    document = Document(path)
    units: list[ExtractedUnit] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        units.append(
            ExtractedUnit(
                text,
                {
                    "source_type": "docx",
                    "paragraph": index,
                    "style": paragraph.style.name if paragraph.style else None,
                },
            )
        )
    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            values = [cell.text.strip().replace("\n", " / ") for cell in row.cells]
            if not any(values):
                continue
            units.append(
                ExtractedUnit(
                    " | ".join(values),
                    {"source_type": "docx", "table": table_index, "row": row_index},
                )
            )
    return units


def extract_pptx(path: Path) -> list[ExtractedUnit]:
    from pptx import Presentation

    presentation = Presentation(path)
    units: list[ExtractedUnit] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text)
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, start=1):
                    values = [cell.text.strip().replace("\n", " / ") for cell in row.cells]
                    if any(values):
                        texts.append(f"表格第{row_index}行：" + " | ".join(values))
        if texts:
            units.append(
                ExtractedUnit(
                    "\n".join(texts),
                    {"source_type": "pptx", "slide": slide_index},
                )
            )
    return units


def _cell_value(value) -> str:  # noqa: ANN001
    if value is None:
        return ""
    return str(value).strip()


def extract_xlsx(path: Path) -> list[ExtractedUnit]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=False)
    units: list[ExtractedUnit] = []
    try:
        for worksheet in workbook.worksheets:
            for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                cells = [_cell_value(value) for value in row]
                if not any(cells):
                    continue
                rendered = " | ".join(f"{index + 1}={value}" for index, value in enumerate(cells) if value)
                units.append(
                    ExtractedUnit(
                        rendered,
                        {"source_type": "xlsx", "sheet": worksheet.title, "row": row_index},
                    )
                )
    finally:
        workbook.close()
    return units


def extract_csv(path: Path, delimiter: str = ",") -> list[ExtractedUnit]:
    text = _read_text(path)
    units: list[ExtractedUnit] = []
    reader = csv.reader(text.splitlines(), delimiter=delimiter)
    for row_index, row in enumerate(reader, start=1):
        values = [value.strip() for value in row]
        if not any(values):
            continue
        rendered = " | ".join(f"{index + 1}={value}" for index, value in enumerate(values) if value)
        units.append(ExtractedUnit(rendered, {"source_type": "csv", "row": row_index}))
    return units


def extract_pdf(path: Path) -> list[ExtractedUnit]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    units: list[ExtractedUnit] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            units.append(ExtractedUnit(text, {"source_type": "pdf", "page": page_index}))
    return units


EXTRACTORS: dict[str, Callable[[Path], list[ExtractedUnit]]] = {
    ".txt": extract_plain,
    ".md": extract_plain,
    ".markdown": extract_plain,
    ".py": extract_plain,
    ".js": extract_plain,
    ".mjs": extract_plain,
    ".ts": extract_plain,
    ".tsx": extract_plain,
    ".json": extract_json,
    ".html": extract_html,
    ".htm": extract_html,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".csv": extract_csv,
    ".tsv": lambda path: extract_csv(path, delimiter="\t"),
    ".pdf": extract_pdf,
}


def supported_suffixes() -> set[str]:
    return set(EXTRACTORS)


def kind_for_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in EXTRACTORS:
        return suffix.lstrip(".")
    return mimetypes.guess_type(path.name)[0] or "unknown"


def extract(path: str | Path) -> list[ExtractedUnit]:
    source = Path(path).resolve()
    extractor = EXTRACTORS.get(source.suffix.lower())
    if not extractor:
        raise ValueError(f"Unsupported knowledge source: {source.suffix or source.name}")
    return extractor(source)
